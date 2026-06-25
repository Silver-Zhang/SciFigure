"""SciFigure minimal Python runtime.

This backend creates an editable PPTX skeleton from a JSON task file.
It is the portable fallback backend for macOS and Linux and can also be
used on Windows when PowerPoint automation is unavailable.
"""

from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover
    raise SystemExit("python-pptx is required. Install dependencies first.") from exc


@dataclass
class TaskSpec:
    name: str
    title: str
    subtitle: str
    core_logic: str
    left_title: str = "Input"
    core_title: str = "Core method"
    right_title: str = "Output"
    output_name: str = "scifigure_output.pptx"


def load_task(path: Path) -> TaskSpec:
    data: dict[str, Any] = json.loads(path.read_text(encoding="utf-8"))
    return TaskSpec(
        name=data.get("name", path.stem),
        title=data.get("title", "SciFigure task"),
        subtitle=data.get("subtitle", "Editable scientific PowerPoint figure"),
        core_logic=data.get("core_logic", "Describe the scientific logic here."),
        left_title=data.get("left_title", "Input"),
        core_title=data.get("core_title", "Core method"),
        right_title=data.get("right_title", "Output"),
        output_name=data.get("output_name", "scifigure_output.pptx"),
    )


def add_textbox(slide, text: str, left, top, width, height, size: int = 16, bold: bool = False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Microsoft YaHei"
    run.font.color.rgb = RGBColor(8, 45, 95) if bold else RGBColor(35, 45, 55)
    return box


def add_panel(slide, title: str, left, top, width, height):
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    panel.line.width = Pt(1.3)
    panel.line.color.rgb = RGBColor(195, 205, 215)
    add_textbox(slide, title, left + Inches(0.15), top + Inches(0.10), width - Inches(0.3), Inches(0.35), 18, True)
    return panel


def add_placeholder(slide, text: str, left, top, width, height):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 247, 250)
    box.line.width = Pt(1)
    box.line.color.rgb = RGBColor(185, 195, 205)
    add_textbox(slide, text, left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), height - Inches(0.16), 12)
    return box


def build_ppt(task: TaskSpec, output_dir: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide, task.title, Inches(0.55), Inches(0.25), Inches(12.2), Inches(0.45), 28, True)
    add_textbox(slide, task.subtitle, Inches(0.55), Inches(0.75), Inches(12.2), Inches(0.30), 15)

    left = Inches(0.55)
    top = Inches(1.30)
    gap = Inches(0.25)
    h = Inches(5.45)
    w_left = Inches(2.25)
    w_core = Inches(7.25)
    w_right = Inches(2.25)

    add_panel(slide, task.left_title, left, top, w_left, h)
    add_panel(slide, task.core_title, left + w_left + gap, top, w_core, h)
    add_panel(slide, task.right_title, left + w_left + gap + w_core + gap, top, w_right, h)

    add_placeholder(slide, "editable input block", left + Inches(0.20), top + Inches(0.80), w_left - Inches(0.40), Inches(1.25))
    add_placeholder(slide, "equation placeholder", left + Inches(0.20), top + Inches(2.25), w_left - Inches(0.40), Inches(0.75))

    add_textbox(slide, task.core_logic, left + w_left + gap + Inches(0.30), top + Inches(0.75), w_core - Inches(0.60), Inches(1.10), 16)
    add_placeholder(slide, "approved local scientific illustration", left + w_left + gap + Inches(0.50), top + Inches(2.05), w_core - Inches(1.00), Inches(2.10))
    add_placeholder(slide, "formula or criterion expression", left + w_left + gap + Inches(0.50), top + Inches(4.45), w_core - Inches(1.00), Inches(0.80))

    add_placeholder(slide, "editable decision block", left + w_left + gap + w_core + gap + Inches(0.20), top + Inches(0.80), w_right - Inches(0.40), Inches(1.20))
    add_placeholder(slide, "editable curve or chart", left + w_left + gap + w_core + gap + Inches(0.20), top + Inches(2.35), w_right - Inches(0.40), Inches(1.75))

    output_dir.mkdir(parents=True, exist_ok=True)
    out = output_dir / task.output_name
    prs.save(out)
    return out


def write_report(task: TaskSpec, output_dir: Path, ppt_path: Path) -> Path:
    report = output_dir / "check_report.md"
    report.write_text(
        "# PPT Quality Check Report\n\n"
        "## Generated file\n\n"
        f"{ppt_path}\n\n"
        "## Status\n\n"
        "Initial editable PPTX skeleton generated by python-pptx backend. Manual style and formula review is still required.\n",
        encoding="utf-8",
    )
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Build an editable SciFigure PPTX skeleton.")
    parser.add_argument("task", type=Path, help="Path to task JSON file")
    parser.add_argument("--output-dir", type=Path, default=Path("result/demo"))
    args = parser.parse_args()

    task = load_task(args.task)
    ppt = build_ppt(task, args.output_dir)
    report = write_report(task, args.output_dir, ppt)
    print(f"[SciFigure] PPTX: {ppt}")
    print(f"[SciFigure] Check report: {report}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
