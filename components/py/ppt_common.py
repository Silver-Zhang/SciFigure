from pptx import Presentation
from pptx.util import Inches


def new_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])
